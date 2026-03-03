"""
generate_slides.py
==================
Reads docs/presentation_spec.md and produces:
  1. Ask_and_You_Shall_Debug_SCaLE.pptx  – the full slide deck
  2. docs/previews/preview_<id>.png       – PNG preview for every slide

Usage (from repo root):
    python slides/generate_slides.py [--no-pptx] [--no-previews]

Requirements:
    pip install python-pptx pillow
"""

import argparse
import os
import textwrap
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Constants – match presentation_spec.md
# ---------------------------------------------------------------------------

PRIMARY = RGBColor(0x11, 0x18, 0x27)   # #111827 near-black
ACCENT  = RGBColor(0x25, 0x63, 0xEB)   # #2563EB blue
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)

FONT_NAME         = "Calibri"
TITLE_PT          = 40
BODY_PT           = 24
FOOTER_PT         = 12
CALLOUT_PT        = 18

FOOTER_TEXT = (
    "Inspektor Gadget MCP Server  •  "
    "https://github.com/inspektor-gadget/ig-mcp-server"
)

# 16:9 Widescreen (13.33″ × 7.5″)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

REPO_ROOT = Path(__file__).parent.parent
SPEC_PATH = REPO_ROOT / "docs" / "presentation_spec.md"
OUT_PPTX  = REPO_ROOT / "Ask_and_You_Shall_Debug_SCaLE.pptx"
PREVIEWS_DIR = REPO_ROOT / "docs" / "previews"

# PIL preview dimensions (1280 × 720 = 16:9)
PX_W, PX_H = 1280, 720

# ---------------------------------------------------------------------------
# Color helpers (hex string → tuples used by PIL)
# ---------------------------------------------------------------------------

def _hex(h: str):
    h = h.lstrip("#")
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

C_PRIMARY    = _hex("111827")
C_ACCENT     = _hex("2563EB")
C_WHITE      = (255, 255, 255)
C_LIGHT_GRAY = (243, 244, 246)
C_MID_GRAY   = (156, 163, 175)

# ---------------------------------------------------------------------------
# Slide data (parsed from presentation_spec.md)
# ---------------------------------------------------------------------------

SLIDES = [
    {
        "id": "title",
        "layout": "title",
        "title": "Ask and You Shall Debug:",
        "subtitle": "Conversational Troubleshooting for Kubernetes",
        "meta": {"right_corner": ["SCaLE"]},
    },
    {
        "id": "why-and-setup",
        "layout": "title_and_bullets",
        "title": "Why conversational troubleshooting?",
        "bullets": [
            "Incidents are time-sensitive and cognitively expensive",
            "Troubleshooting often requires many tools + exact commands",
            "Goal: make real-time observability accessible via natural language",
            "We'll follow an on-call story and resolve multiple scenarios",
        ],
    },
    {
        "id": "agenda",
        "layout": "title_and_bullets",
        "title": "Agenda",
        "bullets": [
            "Demo 1: DNS latency (trace_dns)",
            "Zoom out: Inspektor Gadget + gadgets",
            "Why MCP server (and how it fits)",
            "Mini-demos: disk pressure, CPU saturation, networking misconfig",
            "Wrap-up + resources",
        ],
    },
    {
        "id": "scenario-1-dns-symptoms",
        "layout": "title_and_bullets",
        "title": "Scenario 1: Slow requests → DNS latency",
        "bullets": [
            "Symptom: intermittent slowness / timeouts",
            "Hypothesis space is large (DNS, network, app, load, etc.)",
            "We start by validating whether DNS resolution is slow",
            "Tool we'll use: trace_dns",
        ],
        "callouts": ["Key idea: measure, don't guess."],
    },
    {
        "id": "scenario-1-dns-trace",
        "layout": "title_and_bullets",
        "title": "Using trace_dns to observe DNS in real time",
        "bullets": [
            "trace_dns traces DNS requests and responses",
            "It captures request/response pairs and computes latency",
            "Use it to identify slow lookups and affected workloads",
            "Outcome: narrow the problem quickly and objectively",
        ],
        "images": [
            {
                "path": "docs/media/placeholder_trace_dns.png",
                "caption": "trace_dns output / table",
                "optional": True,
            }
        ],
    },
    {
        "id": "zoom-out-what-is-ig",
        "layout": "title_and_bullets",
        "title": "What is Inspektor Gadget?",
        "bullets": [
            "A systems inspection + data collection framework powered by eBPF",
            "Provides observability in Kubernetes and Linux contexts",
            "Ships a wide selection of \u201cgadgets\u201d for debugging and monitoring",
            "Designed to make low-level visibility usable in familiar workflows",
        ],
    },
    {
        "id": "what-is-a-gadget",
        "layout": "title_and_bullets",
        "title": "What is a \u201cgadget\u201d?",
        "bullets": [
            "A purpose-built observability tool (trace / snapshot / top / profile)",
            "Examples we'll use today:",
            "\u2022 trace_dns (DNS requests/responses + latency)",
            "\u2022 profile_blockio + top_blockio (disk pressure / block I/O)",
            "\u2022 top_process (CPU-heavy processes)",
            "\u2022 tcpdump (packet capture with filters in container context)",
        ],
    },
    {
        "id": "why-mcp",
        "layout": "title_and_bullets",
        "title": "Problem: many gadgets \u2192 hard to choose",
        "bullets": [
            "Inspektor Gadget has a large library of powerful gadgets",
            "In an incident, picking the *right* gadget quickly is the hard part",
            "Teams end up memorizing commands, or reaching for familiar tools first",
            "We wanted: \u201cdescribe the symptom \u2192 get the right tool + next step\u201d",
        ],
        "callouts": ["This is where we discovered MCP."],
    },
    {
        "id": "what-is-mcp",
        "layout": "title_and_bullets",
        "title": "What is MCP (Model Context Protocol)?",
        "bullets": [
            "A standard way for an LLM to call external tools in a structured, safe way",
            "Tools are exposed with names + schemas (inputs/outputs are predictable)",
            "Lets the model choose actions (run tools) and then explain results",
            "In our case: the tools are Inspektor Gadget gadgets via ig-mcp-server",
        ],
        "callouts": ["MCP turns \u201cchat\u201d into \u201cchat + actions\u201d (tool calls)."],
    },
    {
        "id": "architecture-high-level",
        "layout": "title_and_bullets",
        "title": "High-level architecture",
        "bullets": [
            "Copilot/LLM \u2194 MCP protocol \u2194 ig-mcp-server \u2194 Inspektor Gadget gadgets",
            "Gadgets provide real-time signals; LLM turns signals into next steps",
            "MCP server transports: stdio / sse / streamable-http",
            "Gadget discovery: Artifact Hub or explicit gadget images",
        ],
        "images": [
            {
                "path": "docs/media/placeholder_architecture.png",
                "caption": "LLM \u2192 MCP \u2192 IG MCP server \u2192 IG gadgets",
                "optional": True,
            }
        ],
    },
    {
        "id": "human-vs-ai",
        "layout": "title_and_bullets",
        "title": "Human vs. AI troubleshooting",
        "bullets": [
            "Humans: hypothesis-driven, sequential, bias-prone under pressure",
            "LLM + real-time tools: broader search, more consistent method",
            "Best results come from collaboration (human judgment + machine speed)",
            "Guardrail: confirm findings with observable evidence",
        ],
    },
    {
        "id": "scenario-2-disk-pressure",
        "layout": "title_and_bullets",
        "title": "Scenario 2: Disk pressure",
        "bullets": [
            "Symptom: disk pressure alerts, slow writes, cascading failures",
            "Identify which workloads drive I/O and what patterns are present",
            "Gadgets:",
            "\u2022 profile_blockio (profiling block I/O behavior)",
            "\u2022 top_blockio (top block I/O contributors)",
        ],
        "images": [
            {
                "path": "docs/media/placeholder_blockio.png",
                "caption": "profile_blockio / top_blockio output",
                "optional": True,
            }
        ],
    },
    {
        "id": "scenario-3-cpu",
        "layout": "title_and_bullets",
        "title": "Scenario 3: CPU saturation",
        "bullets": [
            "Symptom: high latency + throttling + noisy neighbor suspicion",
            "Goal: identify CPU-heavy processes in the relevant workload context",
            "Gadget: top_process",
            "Optional deeper profiling exists: profile_cpu",
        ],
    },
    {
        "id": "scenario-4-networking",
        "layout": "title_and_bullets",
        "title": "Scenario 4: Networking config/port mismatch",
        "bullets": [
            "Symptom: app can\u2019t connect / sporadic failures",
            "Validate assumptions by observing real packets",
            "Gadget: tcpdump (packet capture with pcap-compatible filters)",
            "Outcome: confirm mismatch between expected vs actual traffic patterns",
        ],
        "images": [
            {
                "path": "docs/media/placeholder_tcpdump.png",
                "caption": "tcpdump gadget output",
                "optional": True,
            }
        ],
    },
    {
        "id": "operations-safety",
        "layout": "title_and_bullets",
        "title": "Operations & safety",
        "bullets": [
            "Run MCP server in read-only mode when appropriate",
            "Use restricted permissions via dedicated service account",
            "Reduce scope to avoid overload (namespace, limits, timeouts)",
            "Treat AI recommendations as hypotheses until validated",
        ],
    },
    {
        "id": "wrap-up",
        "layout": "title_and_bullets",
        "title": "Wrap-up",
        "bullets": [
            "Inspektor Gadget: real-time observability powered by eBPF",
            "MCP server: makes gadgets accessible through an AI interface",
            "Scenarios: DNS latency, disk pressure, CPU saturation, networking",
            "Resources:",
            "\u2022 IG MCP Server: https://github.com/inspektor-gadget/ig-mcp-server",
            "\u2022 Inspektor Gadget: https://github.com/inspektor-gadget/inspektor-gadget",
            "\u2022 Website: https://www.inspektor-gadget.io/",
        ],
    },
]

# ---------------------------------------------------------------------------
# PPTX helpers
# ---------------------------------------------------------------------------

def _rgb(r, g, b):
    return RGBColor(r, g, b)


def _add_text_box(slide, left, top, width, height, text, size_pt,
                  bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox


def _add_rect(slide, left, top, width, height, fill_rgb):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape


def _set_slide_background(slide, rgb: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _add_footer(slide):
    _add_text_box(
        slide,
        left=Inches(0.25), top=SLIDE_H - Inches(0.45),
        width=SLIDE_W - Inches(0.5), height=Inches(0.35),
        text=FOOTER_TEXT,
        size_pt=FOOTER_PT,
        color=C_MID_GRAY_RGB,
        align=PP_ALIGN.CENTER,
    )


C_MID_GRAY_RGB = _rgb(156, 163, 175)


def _add_accent_bar(slide):
    """Thin blue horizontal rule below the title."""
    _add_rect(slide, Inches(0.5), Inches(1.55), SLIDE_W - Inches(1.0), Inches(0.04), ACCENT)


# ---------------------------------------------------------------------------
# PPTX slide builders
# ---------------------------------------------------------------------------

def _build_title_slide(prs):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_background(slide, PRIMARY)

    # Accent bar (left edge decoration)
    _add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, ACCENT)

    # Main title
    _add_text_box(
        slide,
        left=Inches(0.6), top=Inches(1.8),
        width=Inches(10.0), height=Inches(1.5),
        text="Ask and You Shall Debug:",
        size_pt=TITLE_PT + 8,
        bold=True,
        color=WHITE,
    )

    # Subtitle
    _add_text_box(
        slide,
        left=Inches(0.6), top=Inches(3.35),
        width=Inches(10.0), height=Inches(1.0),
        text="Conversational Troubleshooting for Kubernetes",
        size_pt=TITLE_PT - 6,
        bold=False,
        color=_rgb(147, 197, 253),  # light blue
    )

    # Event tag
    _add_text_box(
        slide,
        left=Inches(0.6), top=Inches(4.5),
        width=Inches(4.0), height=Inches(0.5),
        text="SCaLE",
        size_pt=BODY_PT,
        bold=False,
        color=_rgb(156, 163, 175),
    )

    # Organization label bottom-right
    _add_text_box(
        slide,
        left=SLIDE_W - Inches(3.5), top=SLIDE_H - Inches(0.6),
        width=Inches(3.2), height=Inches(0.4),
        text="Inspektor Gadget",
        size_pt=FOOTER_PT,
        bold=True,
        color=_rgb(156, 163, 175),
        align=PP_ALIGN.RIGHT,
    )
    return slide


def _build_content_slide(prs, slide_data):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_background(slide, _rgb(255, 255, 255))

    has_image = bool(slide_data.get("images"))
    content_right = SLIDE_W - Inches(3.8) if has_image else SLIDE_W - Inches(0.5)

    # Title area background strip
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.5), PRIMARY)

    # Title text
    _add_text_box(
        slide,
        left=Inches(0.4), top=Inches(0.2),
        width=content_right - Inches(0.4), height=Inches(1.1),
        text=slide_data["title"],
        size_pt=TITLE_PT - 4,
        bold=True,
        color=WHITE,
    )

    # Accent line
    _add_accent_bar(slide)

    # Bullets
    bullets = slide_data.get("bullets", [])
    callouts = slide_data.get("callouts", [])

    bullet_top = Inches(1.75)
    line_h = Inches(0.48)
    for i, bullet in enumerate(bullets):
        is_sub = bullet.startswith("\u2022")
        indent = Inches(0.55) if is_sub else Inches(0.25)
        _add_text_box(
            slide,
            left=Inches(0.5) + indent, top=bullet_top + i * line_h,
            width=content_right - Inches(0.6) - indent,
            height=line_h,
            text=bullet,
            size_pt=BODY_PT - 4 if is_sub else BODY_PT - 2,
            color=PRIMARY,
        )

    # Callout boxes
    if callouts:
        callout_top = bullet_top + len(bullets) * line_h + Inches(0.15)
        for callout in callouts:
            box = _add_rect(
                slide,
                left=Inches(0.5), top=callout_top,
                width=content_right - Inches(1.0), height=Inches(0.55),
                fill_rgb=_rgb(239, 246, 255),  # very light blue
            )
            _add_text_box(
                slide,
                left=Inches(0.65), top=callout_top + Inches(0.05),
                width=content_right - Inches(1.2), height=Inches(0.48),
                text=callout,
                size_pt=CALLOUT_PT,
                bold=True,
                italic=True,
                color=ACCENT,
            )

    # Placeholder image panel (right side)
    if has_image:
        img_info = slide_data["images"][0]
        img_path = REPO_ROOT / img_info["path"]
        panel_left  = SLIDE_W - Inches(3.6)
        panel_top   = Inches(1.65)
        panel_w     = Inches(3.2)
        panel_h     = Inches(4.4)

        # Draw panel background
        _add_rect(slide, panel_left, panel_top, panel_w, panel_h, LIGHT_GRAY)

        # Try to insert the image; fall back to a labelled placeholder box
        if img_path.exists():
            try:
                slide.shapes.add_picture(
                    str(img_path), panel_left + Inches(0.1), panel_top + Inches(0.1),
                    panel_w - Inches(0.2), panel_h - Inches(0.6),
                )
            except Exception:
                pass

        caption = img_info.get("caption", "")
        if caption:
            _add_text_box(
                slide,
                left=panel_left, top=panel_top + panel_h - Inches(0.55),
                width=panel_w, height=Inches(0.45),
                text=caption,
                size_pt=10,
                color=_rgb(107, 114, 128),
                align=PP_ALIGN.CENTER,
            )

    # Footer
    _add_footer(slide)
    return slide


# ---------------------------------------------------------------------------
# Full PPTX generation
# ---------------------------------------------------------------------------

def generate_pptx(output_path: Path):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    for slide_data in SLIDES:
        if slide_data["layout"] == "title":
            _build_title_slide(prs)
        else:
            _build_content_slide(prs, slide_data)

    prs.save(str(output_path))
    print(f"[pptx] Saved → {output_path}")


# ---------------------------------------------------------------------------
# PIL preview helpers
# ---------------------------------------------------------------------------

def _pil_color(rgb_tuple, alpha=255):
    return (*rgb_tuple, alpha)


def _load_font(size):
    """Try to load Calibri or fall back to a default font."""
    candidates = [
        # Linux – msttcorefonts package
        "/usr/share/fonts/truetype/msttcorefonts/Calibri.ttf",
        # Linux – common fallbacks
        "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
        # macOS
        "/Library/Fonts/Calibri.ttf",
        "/System/Library/Fonts/Helvetica.ttc",
        # Windows
        "C:/Windows/Fonts/calibri.ttf",
        "C:/Windows/Fonts/arial.ttf",
    ]
    for path in candidates:
        if os.path.exists(path):
            try:
                return ImageFont.truetype(path, size)
            except OSError:
                pass
    return ImageFont.load_default()


def _load_font_bold(size):
    candidates = [
        # Linux – msttcorefonts package
        "/usr/share/fonts/truetype/msttcorefonts/Calibrib.ttf",
        # Linux – common fallbacks
        "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
        "/usr/share/fonts/truetype/freefont/FreeSansBold.ttf",
        # macOS
        "/Library/Fonts/CalibriBold.ttf",
        "/System/Library/Fonts/Helvetica.ttc",
        # Windows
        "C:/Windows/Fonts/calibrib.ttf",
        "C:/Windows/Fonts/arialbd.ttf",
    ]
    for path in candidates:
        if os.path.exists(path):
            try:
                return ImageFont.truetype(path, size)
            except OSError:
                pass
    return ImageFont.load_default()


def _draw_wrapped_text(draw, text, font, x, y, max_width, fill, line_height=None):
    """Draw text wrapped to max_width. Returns final y position."""
    words = text.split()
    lines = []
    current = ""
    for word in words:
        test = (current + " " + word).strip()
        bbox = draw.textbbox((0, 0), test, font=font)
        if bbox[2] - bbox[0] <= max_width:
            current = test
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)

    lh = line_height or (draw.textbbox((0, 0), "Ag", font=font)[3] + 4)
    for line in lines:
        draw.text((x, y), line, font=font, fill=fill)
        y += lh
    return y


# ---------------------------------------------------------------------------
# PIL slide renderers
# ---------------------------------------------------------------------------

def _render_title_slide(slide_data) -> Image.Image:
    img = Image.new("RGB", (PX_W, PX_H), C_PRIMARY)
    draw = ImageDraw.Draw(img)

    # Left accent bar
    draw.rectangle([(0, 0), (14, PX_H)], fill=C_ACCENT)

    # Gradient-like overlay strip (subtle)
    for i in range(300):
        alpha = int(30 * (1 - i / 300))
        draw.line([(14, PX_H // 2 - 150 + i), (PX_W, PX_H // 2 - 150 + i)],
                  fill=(*C_ACCENT, alpha))

    title_font   = _load_font_bold(72)
    sub_font     = _load_font(42)
    label_font   = _load_font(24)
    small_font   = _load_font(18)

    # Title
    title = slide_data.get("title", "")
    _draw_wrapped_text(draw, title, title_font, 60, 200, PX_W - 140, C_WHITE, line_height=84)

    # Subtitle
    subtitle = slide_data.get("subtitle", "")
    _draw_wrapped_text(draw, subtitle, sub_font, 60, 340, PX_W - 140, (147, 197, 253), line_height=52)

    # Event tag
    meta = slide_data.get("meta", {})
    for tag in meta.get("right_corner", []):
        draw.text((60, 460), tag, font=label_font, fill=C_MID_GRAY)

    # Org label bottom right
    draw.text((PX_W - 240, PX_H - 50), "Inspektor Gadget", font=small_font, fill=C_MID_GRAY)

    return img


def _render_content_slide(slide_data) -> Image.Image:
    img = Image.new("RGB", (PX_W, PX_H), C_WHITE)
    draw = ImageDraw.Draw(img)

    has_image = bool(slide_data.get("images"))
    content_right = PX_W - 340 if has_image else PX_W - 40

    # Title strip
    TITLE_STRIP_H = 110
    draw.rectangle([(0, 0), (PX_W, TITLE_STRIP_H)], fill=C_PRIMARY)

    title_font   = _load_font_bold(44)
    body_font    = _load_font(26)
    sub_font     = _load_font(22)
    callout_font = _load_font_bold(22)
    footer_font  = _load_font(14)
    caption_font = _load_font(16)

    # Title text
    _draw_wrapped_text(
        draw, slide_data.get("title", ""), title_font,
        30, 18, content_right - 40, C_WHITE, line_height=52,
    )

    # Accent bar
    draw.rectangle([(40, TITLE_STRIP_H + 4), (content_right - 20, TITLE_STRIP_H + 8)],
                   fill=C_ACCENT)

    # Bullets
    bullets = slide_data.get("bullets", [])
    callouts = slide_data.get("callouts", [])

    y = TITLE_STRIP_H + 24
    for bullet in bullets:
        is_sub = bullet.startswith("\u2022")
        bx = 70 if is_sub else 45
        bfont = sub_font if is_sub else body_font
        bcolor = (75, 85, 99) if is_sub else C_PRIMARY

        if not is_sub:
            # Bullet dot
            draw.ellipse([(bx - 16, y + 10), (bx - 8, y + 18)], fill=C_ACCENT)

        y = _draw_wrapped_text(draw, bullet, bfont, bx, y, content_right - bx - 20,
                               bcolor, line_height=36)
        y += 6

    # Callout box
    if callouts:
        y += 10
        for callout in callouts:
            box_h = 48
            draw.rectangle([(40, y), (content_right - 20, y + box_h)],
                           fill=(239, 246, 255))
            draw.rectangle([(40, y), (44, y + box_h)], fill=C_ACCENT)
            draw.text((56, y + 12), callout, font=callout_font, fill=C_ACCENT)
            y += box_h + 10

    # Image panel (right side)
    if has_image:
        panel_x = PX_W - 325
        panel_y = TITLE_STRIP_H + 14
        panel_w = 295
        panel_h = int(PX_H * 0.67)
        draw.rectangle([(panel_x, panel_y), (panel_x + panel_w, panel_y + panel_h)],
                       fill=C_LIGHT_GRAY)
        draw.rectangle([(panel_x, panel_y), (panel_x + panel_w, panel_y + panel_h)],
                       outline=C_ACCENT, width=2)

        # Placeholder icon (simple dashed grid)
        cx, cy = panel_x + panel_w // 2, panel_y + panel_h // 2 - 20
        draw.line([(panel_x + 20, cy), (panel_x + panel_w - 20, cy)],
                  fill=C_MID_GRAY, width=1)
        draw.line([(cx, panel_y + 20), (cx, panel_y + panel_h - 40)],
                  fill=C_MID_GRAY, width=1)

        # Caption
        img_info = slide_data["images"][0]
        caption  = img_info.get("caption", "")
        if caption:
            draw.text((panel_x + 6, panel_y + panel_h + 4), caption,
                      font=caption_font, fill=(107, 114, 128))

    # Footer separator line
    draw.line([(30, PX_H - 36), (PX_W - 30, PX_H - 36)], fill=(229, 231, 235), width=1)

    # Footer text
    draw.text((PX_W // 2, PX_H - 26), FOOTER_TEXT, font=footer_font,
              fill=C_MID_GRAY, anchor="mm")

    return img


def render_preview(slide_data) -> Image.Image:
    if slide_data["layout"] == "title":
        return _render_title_slide(slide_data)
    return _render_content_slide(slide_data)


# ---------------------------------------------------------------------------
# CLI entry point
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Generate PPTX slide deck and/or PNG preview images."
    )
    parser.add_argument("--no-pptx",     action="store_true", help="Skip PPTX generation.")
    parser.add_argument("--no-previews", action="store_true", help="Skip PNG preview generation.")
    parser.add_argument(
        "--preview-slides",
        nargs="*",
        default=None,
        metavar="ID",
        help=(
            "Slide IDs to render as previews. "
            "Defaults to: title agenda scenario-1-dns-symptoms architecture-high-level"
        ),
    )
    args = parser.parse_args()

    default_preview_ids = {
        "title",
        "agenda",
        "scenario-1-dns-symptoms",
        "architecture-high-level",
    }
    preview_ids = set(args.preview_slides) if args.preview_slides is not None else default_preview_ids

    if not args.no_pptx:
        generate_pptx(OUT_PPTX)

    if not args.no_previews:
        PREVIEWS_DIR.mkdir(parents=True, exist_ok=True)
        by_id = {s["id"]: s for s in SLIDES}
        for slide_id in sorted(preview_ids):
            if slide_id not in by_id:
                print(f"[preview] Unknown slide id '{slide_id}' – skipping.")
                continue
            slide_data = by_id[slide_id]
            preview_img = render_preview(slide_data)
            out_path = PREVIEWS_DIR / f"preview_{slide_id}.png"
            preview_img.save(str(out_path))
            print(f"[preview] Saved → {out_path}")


if __name__ == "__main__":
    main()
