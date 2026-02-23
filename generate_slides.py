#!/usr/bin/env python3
"""
Slide generator for "Ask and You Shall Debug" SCaLE presentation.
Reads docs/presentation_spec.md and produces a .pptx file.

Logo (docs/media/inspektor-gadget-logo.png) is placed in the bottom-left
corner of every slide for consistent branding.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Constants derived from the deck spec
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

PRIMARY_COLOR = RGBColor(0x11, 0x18, 0x27)   # #111827 near-black
ACCENT_COLOR = RGBColor(0x25, 0x63, 0xEB)    # #2563EB blue
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)

TITLE_FONT_PT = 40
BODY_FONT_PT = 20
FOOTER_FONT_PT = 12
FONT_NAME = "Calibri"

FOOTER_TEXT = (
    "Inspektor Gadget MCP Server • https://github.com/inspektor-gadget/ig-mcp-server"
)

LOGO_PATH = os.path.join(os.path.dirname(__file__), "docs", "media", "inspektor-gadget-logo.png")
OUTPUT_PATH = os.path.join(os.path.dirname(__file__), "Ask_and_You_Shall_Debug_SCaLE.pptx")

# Logo placement: bottom-left corner of every slide
LOGO_W = Inches(1.6)
LOGO_H = Inches(0.48)
LOGO_LEFT = Inches(0.2)
LOGO_BOTTOM_MARGIN = Inches(0.15)


def _logo_top(slide_h: Emu) -> Emu:
    return slide_h - LOGO_H - LOGO_BOTTOM_MARGIN


def _set_bg(slide, color: RGBColor) -> None:
    """Fill the slide background with a solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_logo(slide) -> None:
    """Add the IG logo to the bottom-left corner of a slide."""
    if not os.path.exists(LOGO_PATH):
        return
    slide.shapes.add_picture(
        LOGO_PATH,
        left=LOGO_LEFT,
        top=_logo_top(SLIDE_H),
        width=LOGO_W,
        height=LOGO_H,
    )


def _add_footer(slide, text: str = FOOTER_TEXT) -> None:
    """Add a small footer text at the bottom centre of a slide."""
    left = Inches(2.0)
    width = SLIDE_W - Inches(4.0)
    top = SLIDE_H - Inches(0.5)
    height = Inches(0.4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(FOOTER_FONT_PT)
    run.font.name = FONT_NAME
    run.font.color.rgb = LIGHT_GRAY


def _add_accent_bar(slide) -> None:
    """Add a thin blue accent bar across the top of a slide."""
    from pptx.util import Inches
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left=Inches(0),
        top=Inches(0),
        width=SLIDE_W,
        height=Inches(0.07),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_COLOR
    bar.line.fill.background()


def _add_title_text(slide, title: str, left, top, width, height,
                    font_pt: int = TITLE_FONT_PT, bold: bool = True,
                    color: RGBColor = WHITE) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.name = FONT_NAME
    run.font.color.rgb = color


def _add_bullets(slide, bullets: list[str], left, top, width, height) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = bullet if bullet.startswith("•") else f"• {bullet}"
        run.font.size = Pt(BODY_FONT_PT)
        run.font.name = FONT_NAME
        run.font.color.rgb = WHITE


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def build_title_slide(prs: Presentation) -> None:
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    _set_bg(slide, PRIMARY_COLOR)
    _add_accent_bar(slide)

    # Main title
    _add_title_text(
        slide,
        "Ask and You Shall Debug:",
        left=Inches(0.8),
        top=Inches(1.5),
        width=Inches(11.0),
        height=Inches(1.4),
        font_pt=44,
    )
    # Subtitle
    _add_title_text(
        slide,
        "Conversational Troubleshooting for Kubernetes",
        left=Inches(0.8),
        top=Inches(3.0),
        width=Inches(11.0),
        height=Inches(1.0),
        font_pt=28,
        bold=False,
        color=RGBColor(0xA0, 0xC4, 0xFF),
    )
    # Event label
    _add_title_text(
        slide,
        "Inspektor Gadget MCP Server  •  SCaLE",
        left=Inches(0.8),
        top=Inches(4.3),
        width=Inches(11.0),
        height=Inches(0.6),
        font_pt=18,
        bold=False,
        color=LIGHT_GRAY,
    )

    _add_logo(slide)


def build_content_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    _set_bg(slide, PRIMARY_COLOR)
    _add_accent_bar(slide)

    # Slide title
    _add_title_text(
        slide,
        title,
        left=Inches(0.6),
        top=Inches(0.3),
        width=Inches(12.0),
        height=Inches(1.1),
        font_pt=TITLE_FONT_PT,
    )

    # Bullets
    _add_bullets(
        slide,
        bullets,
        left=Inches(0.6),
        top=Inches(1.55),
        width=Inches(12.0),
        height=Inches(4.8),
    )

    _add_footer(slide)
    _add_logo(slide)


# ---------------------------------------------------------------------------
# Slide definitions from the spec
# ---------------------------------------------------------------------------

SLIDES = [
    # Slide 1 – title (handled separately)
    None,
    # Slide 2
    {
        "title": "Why conversational troubleshooting?",
        "bullets": [
            "Incidents are time-sensitive and cognitively expensive",
            "Troubleshooting often requires many tools + exact commands",
            "Goal: make real-time observability accessible via natural language",
            "We'll follow an on-call story and resolve multiple scenarios",
        ],
    },
    # Slide 3 – agenda
    {
        "title": "Agenda",
        "bullets": [
            "Demo 1: DNS latency (trace_dns)",
            "Zoom out: Inspektor Gadget + gadgets",
            "Why MCP server (and how it fits)",
            "Mini-demos: disk pressure, CPU saturation, networking misconfig",
            "Wrap-up + resources",
        ],
    },
    # Slide 4 – scenario 1 symptoms
    {
        "title": "Scenario 1: Slow requests → DNS latency",
        "bullets": [
            "Symptom: intermittent slowness / timeouts",
            "Hypothesis space is large (DNS, network, app, load, etc.)",
            "We start by validating whether DNS resolution is slow",
            "Tool we'll use: trace_dns",
            "Key idea: measure, don't guess.",
        ],
    },
    # Slide 5
    {
        "title": "Using trace_dns to observe DNS in real time",
        "bullets": [
            "trace_dns traces DNS requests and responses",
            "It captures request/response pairs and computes latency",
            "Use it to identify slow lookups and affected workloads",
            "Outcome: narrow the problem quickly and objectively",
        ],
    },
    # Slide 6
    {
        "title": "What is Inspektor Gadget?",
        "bullets": [
            "A systems inspection + data collection framework powered by eBPF",
            "Provides observability in Kubernetes and Linux contexts",
            'Ships a wide selection of \u201cgadgets\u201d for debugging and monitoring',
            "Designed to make low-level visibility usable in familiar workflows",
        ],
    },
    # Slide 7
    {
        "title": 'What is a \u201cgadget\u201d?',
        "bullets": [
            "A purpose-built observability tool (trace / snapshot / top / profile)",
            "Examples we'll use today:",
            "• trace_dns (DNS requests/responses + latency)",
            "• profile_blockio + top_blockio (disk pressure / block I/O)",
            "• top_process (CPU-heavy processes)",
            "• tcpdump (packet capture with filters in container context)",
        ],
    },
    # Slide 8
    {
        "title": "Problem: many gadgets → hard to choose",
        "bullets": [
            "Inspektor Gadget has a large library of powerful gadgets",
            "In an incident, picking the right gadget quickly is the hard part",
            "Teams end up memorizing commands, or reaching for familiar tools first",
            "We wanted: \u201cdescribe the symptom \u2192 get the right tool + next step\u201d",
            "This is where we discovered MCP.",
        ],
    },
    # Slide 9
    {
        "title": "What is MCP (Model Context Protocol)?",
        "bullets": [
            "A standard way for an LLM to call external tools in a structured, safe way",
            "Tools are exposed with names + schemas (inputs/outputs are predictable)",
            "Lets the model choose actions (run tools) and then explain results",
            "In our case: the tools are Inspektor Gadget gadgets via ig-mcp-server",
            'MCP turns \u201cchat\u201d into \u201cchat + actions\u201d (tool calls).',
        ],
    },
    # Slide 10 – architecture
    {
        "title": "High-level architecture",
        "bullets": [
            "Copilot/LLM ↔ MCP protocol ↔ ig-mcp-server ↔ Inspektor Gadget gadgets",
            "Gadgets provide real-time signals; LLM turns signals into next steps",
            "MCP server transports: stdio / sse / streamable-http",
            "Gadget discovery: Artifact Hub or explicit gadget images",
            "[See: docs/media/placeholder_architecture.png]",
        ],
    },
    # Slide 11
    {
        "title": "Human vs. AI troubleshooting",
        "bullets": [
            "Humans: hypothesis-driven, sequential, bias-prone under pressure",
            "LLM + real-time tools: broader search, more consistent method",
            "Best results come from collaboration (human judgment + machine speed)",
            "Guardrail: confirm findings with observable evidence",
        ],
    },
    # Slide 12
    {
        "title": "Scenario 2: Disk pressure",
        "bullets": [
            "Symptom: disk pressure alerts, slow writes, cascading failures",
            "Identify which workloads drive I/O and what patterns are present",
            "Gadgets:",
            "• profile_blockio (profiling block I/O behavior)",
            "• top_blockio (top block I/O contributors)",
        ],
    },
    # Slide 13
    {
        "title": "Scenario 3: CPU saturation",
        "bullets": [
            "Symptom: high latency + throttling + noisy neighbor suspicion",
            "Goal: identify CPU-heavy processes in the relevant workload context",
            "Gadget: top_process",
            "Optional deeper profiling exists: profile_cpu",
        ],
    },
    # Slide 14
    {
        "title": "Scenario 4: Networking config/port mismatch",
        "bullets": [
            "Symptom: app can't connect / sporadic failures",
            "Validate assumptions by observing real packets",
            "Gadget: tcpdump (packet capture with pcap-compatible filters)",
            "Outcome: confirm mismatch between expected vs actual traffic patterns",
        ],
    },
    # Slide 15
    {
        "title": "Operations & safety",
        "bullets": [
            "Run MCP server in read-only mode when appropriate",
            "Use restricted permissions via dedicated service account",
            "Reduce scope to avoid overload (namespace, limits, timeouts)",
            "Treat AI recommendations as hypotheses until validated",
        ],
    },
    # Slide 16
    {
        "title": "Wrap-up",
        "bullets": [
            "Inspektor Gadget: real-time observability powered by eBPF",
            "MCP server: makes gadgets accessible through an AI interface",
            "Scenarios: DNS latency, disk pressure, CPU saturation, networking",
            "Resources:",
            "• IG MCP Server: https://github.com/inspektor-gadget/ig-mcp-server",
            "• Inspektor Gadget: https://github.com/inspektor-gadget/inspektor-gadget",
            "• Website: https://www.inspektor-gadget.io/",
        ],
    },
]


# ---------------------------------------------------------------------------
# Preview image generator (Pillow)
# ---------------------------------------------------------------------------

def generate_preview(title: str, bullets: list[str], output_path: str,
                     is_title_slide: bool = False) -> None:
    """Render a simple PNG preview of a slide using Pillow."""
    from PIL import Image, ImageDraw, ImageFont

    W, H = 1333, 750
    img = Image.new("RGB", (W, H), (0x11, 0x18, 0x27))
    draw = ImageDraw.Draw(img)

    # Accent bar at top
    draw.rectangle([0, 0, W, 7], fill=(0x25, 0x63, 0xEB))

    # Load fonts
    try:
        font_title = ImageFont.truetype(
            "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 42 if not is_title_slide else 52
        )
        font_subtitle = ImageFont.truetype(
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 30
        )
        font_body = ImageFont.truetype(
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 22
        )
        font_footer = ImageFont.truetype(
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 13
        )
    except Exception:
        font_title = ImageFont.load_default()
        font_subtitle = font_body = font_footer = font_title

    if is_title_slide:
        # Title slide layout
        draw.text((80, 140), title, fill=(0xFF, 0xFF, 0xFF), font=font_title)
        draw.text((80, 215), bullets[0] if bullets else "", fill=(0xA0, 0xC4, 0xFF), font=font_subtitle)
        draw.text((80, 400), "Inspektor Gadget MCP Server  •  SCaLE",
                  fill=(0xCC, 0xCC, 0xCC), font=font_body)
    else:
        # Content slide layout
        draw.text((60, 30), title, fill=(0xFF, 0xFF, 0xFF), font=font_title)
        # Divider line under title
        draw.rectangle([60, 100, W - 60, 103], fill=(0x25, 0x63, 0xEB))

        y = 120
        for bullet in bullets:
            text = bullet if bullet.startswith("•") else f"• {bullet}"
            draw.text((80, y), text, fill=(0xFF, 0xFF, 0xFF), font=font_body)
            y += 48

        # Footer
        draw.text(
            (W // 2 - 300, H - 42),
            FOOTER_TEXT,
            fill=(0xCC, 0xCC, 0xCC),
            font=font_footer,
        )

    # Logo in bottom-left corner
    logo_path = LOGO_PATH
    if os.path.exists(logo_path):
        try:
            logo = Image.open(logo_path).convert("RGBA")
            logo_w = 160
            logo_h = int(logo.height * logo_w / logo.width)
            logo = logo.resize((logo_w, logo_h), Image.Resampling.LANCZOS)
            logo_x = 20
            logo_y = H - logo_h - 15
            img.paste(logo, (logo_x, logo_y), logo)
        except Exception as e:
            print(f"Warning: could not paste logo: {e}")

    img.save(output_path)
    print(f"Preview saved: {output_path}")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Slide 1 – title
    build_title_slide(prs)

    # Slides 2-16 – content
    for slide_def in SLIDES[1:]:
        build_content_slide(prs, slide_def["title"], slide_def["bullets"])

    prs.save(OUTPUT_PATH)
    print(f"Presentation saved: {OUTPUT_PATH}")

    # Generate preview images
    media_dir = os.path.join(os.path.dirname(__file__), "docs", "media")
    os.makedirs(media_dir, exist_ok=True)

    # Title slide preview
    generate_preview(
        "Ask and You Shall Debug:",
        ["Conversational Troubleshooting for Kubernetes"],
        os.path.join(media_dir, "preview_title.png"),
        is_title_slide=True,
    )

    # Agenda slide preview (slide 3, index 2 in SLIDES)
    agenda = SLIDES[2]
    generate_preview(
        agenda["title"],
        agenda["bullets"],
        os.path.join(media_dir, "preview_agenda.png"),
    )

    # Scenario slide preview (slide 4, index 3 in SLIDES)
    scenario = SLIDES[3]
    generate_preview(
        scenario["title"],
        scenario["bullets"],
        os.path.join(media_dir, "preview_scenario.png"),
    )

    # Architecture slide preview (slide 10, index 9 in SLIDES)
    arch = SLIDES[9]
    generate_preview(
        arch["title"],
        arch["bullets"],
        os.path.join(media_dir, "preview_architecture.png"),
    )


if __name__ == "__main__":
    main()
